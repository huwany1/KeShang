"""
模块: services.document_service.worker
职责: 文档处理 Celery Worker：下载 -> 解析 -> NER/RE -> 入库/入图 -> 状态与指标。
输入: object_path。
输出: 解析统计。
"""

from pathlib import Path
from typing import List, Tuple

from celery import shared_task
from sqlalchemy import update
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.orm import sessionmaker
from sqlalchemy.ext.asyncio import create_async_engine

import fitz  # PyMuPDF
from pptx import Presentation  # python-pptx

from common.mq.celery_app import celery_app
from common.config.settings import settings
from common.storage.minio_client import get_minio_client, ensure_bucket, get_object_to_path
from common.graph.neo4j_writer import ensure_document_and_concepts, create_related_edges
from common.cache.redis_client import get_redis_client
from services.document_service.models import Document, DocumentConcept, DocumentRelation


@shared_task(name="document.extract_text", autoretry_for=(Exception,), retry_kwargs={"max_retries": 3, "countdown": 5})
def extract_text_task(object_path: str) -> dict:
    """解析对象存储中的文件，使用可插拔 NER/RE 流水线抽取概念与关系，入库/入图并更新状态。
    输入: object_path 对象存储键名。
    输出: {document_id, num_concepts}。
    作用: 支撑 MVP 与后续替换为更强的管线（LLM/专用模型）。
    """
    metrics = get_redis_client(db=2)
    try:
        client = get_minio_client()
        ensure_bucket(client)
        tmp_dir = Path(".build/worker")
        tmp_dir.mkdir(parents=True, exist_ok=True)
        local_path = tmp_dir / object_path.replace("/", "_")
        get_object_to_path(client, object_path, local_path)

        text: str = _dispatch_parse(local_path)
        entities, relations = _pipeline_ner_re(text)

        document_id = _derive_doc_id_from_object(object_path)
        _persist_granular_results(document_id, entities, relations)
        ensure_document_and_concepts(document_id, entities)
        if relations:
            create_related_edges(relations)

        _update_document_status(document_id, status="ready")
        metrics.incr("metrics:document_processed")
        return {"document_id": document_id, "num_concepts": len(entities)}
    except Exception as exc:
        # 标记失败并计数；Celery 会按装饰器策略重试
        try:
            document_id = _derive_doc_id_from_object(object_path)
            _update_document_status(document_id, status="failed")
        except Exception:
            pass
        metrics.incr("metrics:document_failed")
        raise exc


def _dispatch_parse(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _parse_pdf(path)
    if suffix in (".ppt", ".pptx"):
        return _parse_ppt(path)
    return ""


def _parse_pdf(path: Path) -> str:
    doc = fitz.open(str(path))
    parts: List[str] = []
    for page in doc:
        parts.append(page.get_text())
    return "\n".join(parts)


def _parse_ppt(path: Path) -> str:
    prs = Presentation(str(path))
    parts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
    return "\n".join(parts)


def _pipeline_ner_re(text: str) -> Tuple[List[str], List[Tuple[str, str]]]:
    """可插拔 NER/RE 流水线（占位）：
    - NER: 规则分词 + 频次/长度过滤
    - RE: 相邻共现形成临时边
    输入: 文本。
    输出: (实体列表, 关系列表)。
    作用: 后续可替换为 spaCy/Transformers/LLM。
    """
    if not text:
        return [], []
    import re
    tokens = re.split(r"[^\w\u4e00-\u9fa5]+", text)
    tokens = [t.strip() for t in tokens if len(t.strip()) >= 2]
    # 频次过滤
    freq = {}
    for t in tokens:
        freq[t] = freq.get(t, 0) + 1
    entities = [t for t, c in freq.items() if c >= 2][:100]
    relations: List[Tuple[str, str]] = []
    for i in range(len(entities) - 1):
        relations.append((entities[i], entities[i + 1]))
    return entities, relations


def _persist_granular_results(document_id: str, entities: List[str], relations: List[Tuple[str, str]]) -> None:
    database_url = (
        f"postgresql+asyncpg://{settings.postgresUser}:{settings.postgresPassword}"
        f"@{settings.postgresHost}:{settings.postgresPort}/{settings.postgresDb}"
    )
    engine = create_async_engine(database_url, future=True, echo=False)
    async_session_factory = sessionmaker(bind=engine, class_=AsyncSession, expire_on_commit=False)

    async def _do_persist() -> None:
        async with async_session_factory() as session:
            for name in entities:
                session.add(DocumentConcept(documentId=document_id, conceptName=name))
            for a, b in relations:
                session.add(DocumentRelation(documentId=document_id, sourceConcept=a, targetConcept=b))
            await session.commit()

    import asyncio
    asyncio.get_event_loop().run_until_complete(_do_persist())


def _update_document_status(document_id: str, status: str) -> None:
    database_url = (
        f"postgresql+asyncpg://{settings.postgresUser}:{settings.postgresPassword}"
        f"@{settings.postgresHost}:{settings.postgresPort}/{settings.postgresDb}"
    )
    engine = create_async_engine(database_url, future=True, echo=False)
    async_session_factory = sessionmaker(bind=engine, class_=AsyncSession, expire_on_commit=False)

    async def _do_update() -> None:
        async with async_session_factory() as session:
            await session.execute(update(Document).where(Document.id == document_id).values(status=status))
            await session.commit()

    import asyncio
    asyncio.get_event_loop().run_until_complete(_do_update())


# 确保模块被导入时 Celery 应用载入
app = celery_app
